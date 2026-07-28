# Copyright Advanced Micro Devices, Inc., or its affiliates.
# SPDX-License-Identifier: MIT
"""Generate the Instruction-Prefetch-Passes deck (Prefetch-Passes-Overview.pptx).

15-slide technical deck: graphic-first (cards, pills, size ladders, tables) with
compact mono boxes for instruction bursts/operands. Keeps the key numbers
(CP boundary 32640 B, spacing 4096 B, I-cache 65536 B / 64 KB).
Run: `python3 generate_prefetch_slides.py`.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

import os

_THEME = os.environ.get("DECK_THEME", "dark").lower()

if _THEME in ("amd", "light"):
    # AMD branding on a dark background: Arial + AMD brand accents.
    BG = RGBColor(0x0A, 0x0A, 0x0A)
    PANEL = RGBColor(0x17, 0x17, 0x1A)
    PANEL2 = RGBColor(0x1F, 0x1F, 0x24)
    DARK = RGBColor(0x12, 0x12, 0x15)
    CODEBG = RGBColor(0x0E, 0x0E, 0x11)
    WHITE = RGBColor(0xF2, 0xF2, 0xF2)
    GREY = RGBColor(0xB4, 0xB4, 0xB4)
    DIM = RGBColor(0x8A, 0x8A, 0x8A)
    RED = RGBColor(0xED, 0x1C, 0x24)  # AMD red
    GREEN = RGBColor(0x3F, 0xC1, 0x74)
    BLUE = RGBColor(0x00, 0xC2, 0xDE)  # AMD teal
    AMBER = RGBColor(0xF2, 0x65, 0x22)  # AMD orange
    LINEC = RGBColor(0x33, 0x33, 0x38)
    INK = RGBColor(0x0A, 0x0A, 0x0A)  # text on saturated color fills
    CELL_LINE = RGBColor(0x2A, 0x2A, 0x30)
    BEST_FILL = RGBColor(0x1B, 0x2A, 0x1B)
    MONO = "Consolas"
    SANS = "Arial"
    ACCENT_BAR = False  # AMD style: no red bars
    OUT = "Prefetch-Passes-Overview-AMD.pptx"
else:
    # Original dark theme.
    BG = RGBColor(0x0A, 0x0A, 0x0A)
    PANEL = RGBColor(0x17, 0x17, 0x1A)
    PANEL2 = RGBColor(0x1F, 0x1F, 0x24)
    DARK = RGBColor(0x12, 0x12, 0x15)
    CODEBG = RGBColor(0x0E, 0x0E, 0x11)
    WHITE = RGBColor(0xF2, 0xF2, 0xF2)
    GREY = RGBColor(0xB4, 0xB4, 0xB4)
    DIM = RGBColor(0x8A, 0x8A, 0x8A)
    RED = RGBColor(0xE6, 0x3B, 0x3B)
    GREEN = RGBColor(0x4C, 0xC2, 0x7A)
    BLUE = RGBColor(0x4A, 0x9E, 0xE0)
    AMBER = RGBColor(0xE0, 0xA1, 0x3A)
    LINEC = RGBColor(0x33, 0x33, 0x38)
    INK = RGBColor(0x0A, 0x0A, 0x0A)
    CELL_LINE = RGBColor(0x2A, 0x2A, 0x30)
    BEST_FILL = RGBColor(0x1B, 0x2A, 0x1B)
    MONO = "Consolas"
    SANS = "Segoe UI"
    ACCENT_BAR = True
    OUT = "Prefetch-Passes-Overview.pptx"

# Kernel the slide-15/16 performance data was measured on (bbs_pgr2_sia0.yaml).
KCFG = (
    "Measured on:  BBS bf16 GEMM (TN)  .  256x256 tile  .  DepthU 128  .  PGR2  .  "
    "M2048 x N1024 x K4096  .  gfx1250  .  I-cache rotating buffer forces the miss"
)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
EMU_W, EMU_H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    return sp


def txt(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    return tf


def para(
    tf,
    runs,
    size=18,
    align=PP_ALIGN.LEFT,
    space_after=6,
    level=0,
    line_spacing=1.06,
    bullet=None,
    first=False,
):
    p = (
        tf.paragraphs[0]
        if (first and not tf.paragraphs[0].runs)
        else tf.add_paragraph()
    )
    p.alignment = align
    p.level = level
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    p.line_spacing = line_spacing
    if isinstance(runs, str):
        runs = [(runs, WHITE, False, SANS)]
    for text, color, bold, font in runs:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    pPr = p._p.get_or_add_pPr()
    if bullet is None:
        pPr.append(pPr.makeelement(qn("a:buNone"), {}))
    else:
        pPr.append(pPr.makeelement(qn("a:buChar"), {"char": bullet}))
    return p


def header(s, kicker, title, page=None):
    if ACCENT_BAR:
        rect(s, Inches(0), Inches(0), Inches(0.16), EMU_H, fill=RED)
    tf = txt(s, Inches(0.55), Inches(0.32), Inches(12.3), Inches(0.3))
    para(tf, [(kicker, RED, True, MONO)], size=12, space_after=0, first=True)
    tf2 = txt(s, Inches(0.55), Inches(0.58), Inches(12.3), Inches(0.75))
    para(tf2, [(title, WHITE, True, SANS)], size=28, space_after=0, first=True)
    if ACCENT_BAR:
        rect(s, Inches(0.58), Inches(1.3), Inches(2.0), Inches(0.035), fill=RED)
    if page is not None:
        tf3 = txt(s, Inches(12.4), Inches(7.04), Inches(0.8), Inches(0.32))
        para(
            tf3,
            [(str(page), DIM, False, MONO)],
            size=11,
            align=PP_ALIGN.RIGHT,
            space_after=0,
            first=True,
        )


def card(s, x, y, w, h, fill=PANEL, line=LINEC, line_w=1.0):
    return rect(
        s,
        x,
        y,
        w,
        h,
        fill=fill,
        line=line,
        line_w=line_w,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE,
    )


def picture(s, path, x, y, w):
    return s.shapes.add_picture(path, x, y, width=w)


def pill(s, x, y, w, name, color=RED, h=0.46, size=12):
    p = rect(
        s,
        x,
        y,
        w,
        Inches(h),
        fill=PANEL2,
        line=color,
        line_w=1.25,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE,
    )
    tf = p.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = tf.margin_bottom = Pt(1)
    para(
        tf,
        [(name, color, True, MONO)],
        size=size,
        align=PP_ALIGN.CENTER,
        space_after=0,
        first=True,
    )
    return p


def codebox(s, x, y, w, h, lines, size=13, title=None):
    card(s, x, y, w, h, fill=CODEBG, line=LINEC)
    tf = txt(s, x + Inches(0.2), y + Inches(0.12), w - Inches(0.4), h - Inches(0.24))
    first = True
    if title:
        para(tf, [(title, DIM, True, MONO)], size=11, space_after=6, first=True)
        first = False
    for ln in lines:
        runs = [ln] if isinstance(ln, tuple) else [(ln, GREY, False, MONO)]
        para(tf, runs, size=size, space_after=2, line_spacing=1.0, first=first)
        first = False
    return tf


def bullets(s, x, y, w, h, items, size=15, gap=10):
    tf = txt(s, x, y, w, h)
    for i, it in enumerate(items):
        runs = it if isinstance(it, list) else [(it, WHITE, False, SANS)]
        para(tf, runs, size=size, space_after=gap, bullet="•", first=(i == 0))
    return tf


def status_line(s, x, y, w, status=None, challenge=None):
    if status:
        t = txt(s, x, y, w, Inches(0.4))
        para(
            t,
            [("Status:  ", GREEN, True, MONO), (status, GREY, False, SANS)],
            size=11,
            space_after=0,
            first=True,
        )
        y = y + Inches(0.4)
    if challenge:
        t = txt(s, x, y, w, Inches(0.5))
        para(
            t,
            [("Challenge:  ", RED, True, MONO), (challenge, GREY, False, SANS)],
            size=11,
            space_after=0,
            first=True,
        )


# =========================================================================
# 1. TITLE
# =========================================================================
s = slide()
if ACCENT_BAR:
    rect(s, Inches(0), Inches(0), EMU_W, Inches(0.16), fill=RED)
    rect(s, Inches(0), EMU_H - Inches(0.16), EMU_W, Inches(0.16), fill=RED)
tf = txt(s, Inches(1.0), Inches(2.4), Inches(11.3), Inches(1.5))
para(
    tf,
    [("Instruction Prefetch", WHITE, True, SANS)],
    size=46,
    space_after=6,
    first=True,
)
para(
    tf,
    [
        (
            "CP & SW instruction prefetch for MI450 / gfx1250 (StinkyTofu)",
            GREY,
            False,
            SANS,
        )
    ],
    size=22,
    space_after=0,
)
tf2 = txt(s, Inches(1.02), Inches(4.35), Inches(11.3), Inches(0.5))
para(
    tf2,
    [
        ("Relative (PC-rel)", BLUE, True, SANS),
        ("  vs  ", DIM, False, SANS),
        ("Absolute (getpc)", AMBER, True, SANS),
        ("     Static / Dynamic", DIM, False, SANS),
    ],
    size=16,
    space_after=0,
    first=True,
)
tf3 = txt(s, Inches(1.02), Inches(5.05), Inches(11.3), Inches(0.4))
para(
    tf3,
    [("George Tseng   |   ML Libraries TW", DIM, False, SANS)],
    size=15,
    space_after=0,
    first=True,
)

# =========================================================================
# 2. AGENDA
# =========================================================================
s = slide()
header(s, "AGENDA", "What we will cover", 2)
left = [
    "Motivation — why Instruction prefetch",
    "Prefetch taxonomy",
    "CP prefetch & instruction sizing",
    "The hardware hint (operands)",
    "Global byte grid & boundaries",
    "Relative: overview",
    "Relative: grid-based insertion",
]
right = [
    "Absolute: overview",
    "Size-based tier selection & pipeline",
    "Absolute static: burst",
    "Absolute dynamic (CFG-target)",
    "Relative vs Absolute",
    "Cache-miss evidence",
    "Results & status",
]
for x, col, start in [(Inches(0.85), left, 1), (Inches(6.95), right, 8)]:
    tf = txt(s, x, Inches(1.8), Inches(5.9), Inches(5.4))
    for i, it in enumerate(col):
        para(
            tf,
            [(f"{start+i:02d}  ", RED, True, MONO), (it, WHITE, False, SANS)],
            size=17,
            space_after=12,
            bullet=None,
            first=(i == 0),
        )

# =========================================================================
# 3. MOTIVATION
# =========================================================================
s = slide()
header(s, "MOTIVATION", "Why software instruction prefetch?", 3)
tf = txt(s, Inches(0.58), Inches(1.7), Inches(6.7), Inches(2.7))
para(
    tf,
    [
        ("PROBLEM   ", RED, True, MONO),
        ("MI450 has ", GREY, False, SANS),
        ("no HW instruction prefetch", WHITE, True, SANS),
        ("; a ", GREY, False, SANS),
        ("128 B", WHITE, True, MONO),
        (" I-cache miss costs ", GREY, False, SANS),
        ("~1000 cycles", RED, True, MONO),
        (" — big kernels stall repeatedly.", GREY, False, SANS),
    ],
    size=15,
    space_after=11,
    first=True,
    bullet=None,
)
para(
    tf,
    [
        ("GOAL   ", AMBER, True, MONO),
        ("Cut repeated I-cache miss latency — keep fetch ", GREY, False, SANS),
        ("ahead of the program counter", WHITE, True, SANS),
        (".", GREY, False, SANS),
    ],
    size=15,
    space_after=11,
    bullet=None,
)
para(
    tf,
    [
        ("SOLUTION   ", GREEN, True, MONO),
        ("CP prefetch (HW, kernel head)  +  ", GREY, False, SANS),
        ("SW prefetch", WHITE, True, SANS),
        ("  (s_prefetch_inst, the tail).", GREY, False, SANS),
    ],
    size=15,
    space_after=0,
    bullet=None,
)
# simple CP-vs-tail coverage bar (left, under the bullets)
_by, _bh, _frac, _bw = Inches(4.55), Inches(0.7), 0.34, 6.4
_b1 = rect(s, Inches(0.58), _by, Inches(_bw * _frac), _bh, fill=GREEN)
_b2 = rect(
    s,
    Inches(0.58) + Inches(_bw * _frac),
    _by,
    Inches(_bw * (1 - _frac)),
    _bh,
    fill=PANEL2,
)
rect(s, Inches(0.58), _by, Inches(_bw), _bh, fill=None, line=LINEC)
_t = _b1.text_frame
_t.vertical_anchor = MSO_ANCHOR.MIDDLE
para(
    _t,
    [("CP", INK, True, SANS)],
    size=14,
    align=PP_ALIGN.CENTER,
    space_after=0,
    first=True,
)
_t = _b2.text_frame
_t.vertical_anchor = MSO_ANCHOR.MIDDLE
_t.word_wrap = True
para(
    _t,
    [("stall risk", GREY, True, SANS)],
    size=14,
    align=PP_ALIGN.CENTER,
    space_after=0,
    first=True,
)
tf = txt(s, Inches(0.58), Inches(5.35), Inches(6.5), Inches(1.2))
para(
    tf,
    [
        ("CP covers ", GREY, False, SANS),
        ("[0, 32640 B)", GREEN, True, MONO),
        (";", GREY, False, SANS),
    ],
    size=16,
    space_after=4,
    first=True,
)
para(
    tf,
    [
        (
            "SW prefetch covers the tail so fetch stays ahead of the program counter.",
            GREY,
            False,
            SANS,
        )
    ],
    size=16,
    space_after=0,
)
picture(s, "evidence/miss_no_prefetch.png", Inches(7.4), Inches(1.85), Inches(5.35))
tf = txt(s, Inches(7.4), Inches(5.28), Inches(5.35), Inches(1.5))
para(
    tf,
    [("Baseline (no prefetch): I-cache miss latency", WHITE, True, SANS)],
    size=14,
    space_after=3,
    first=True,
)
para(
    tf,
    [
        ("mean miss latency ", GREY, False, SANS),
        ("~400-600 cycles", RED, True, MONO),
        (" (single misses peak ~1000) — what prefetch removes.", GREY, False, SANS),
    ],
    size=13,
    space_after=0,
)

# =========================================================================
# 4. HARDWARE HINT
# =========================================================================
s = slide()
header(s, "ISA", "The hardware hint: s_prefetch_inst", 6)
codebox(
    s,
    Inches(0.58),
    Inches(1.7),
    Inches(5.9),
    Inches(1.5),
    title="ABSOLUTE",
    lines=[
        ("s_prefetch_inst s[base:base+1],", AMBER, True, MONO),
        ("                koffset, null, 0x1f", AMBER, True, MONO),
        ("base pair = absolute address", GREY, False, SANS),
    ],
    size=13,
)
codebox(
    s,
    Inches(6.6),
    Inches(1.7),
    Inches(6.15),
    Inches(1.5),
    title="PC-RELATIVE",
    lines=[
        ("s_prefetch_inst_pc_rel 0, null, 31", BLUE, True, MONO),
        ("", GREY, False, MONO),
        ("address = PC + 8 + koffset  (no SGPR)", GREY, False, SANS),
    ],
    size=13,
)
bullets(
    s,
    Inches(0.58),
    Inches(3.5),
    Inches(12.2),
    Inches(3.3),
    [
        [
            ("koffset", WHITE, True, MONO),
            (
                "  — step within the target window (abs) / PC offset (pc-rel).",
                GREY,
                False,
                SANS,
            ),
        ],
        [
            ("slength = null", WHITE, True, MONO),
            ("  — no SGPR supplies the length (= 0).", GREY, False, SANS),
        ],
        [
            ("klength = 31 (0x1f)", WHITE, True, MONO),
            ("  — length = ((0)+(31)) & 31 + 1 = ", GREY, False, SANS),
            ("32 cache lines", WHITE, True, SANS),
            (".", GREY, False, SANS),
        ],
        [
            ("Fetch length = 32 lines x 128 B = ", GREY, False, SANS),
            ("4096 B", GREEN, True, MONO),
            (" per hint (klength tunable up to ~32 KB).", GREY, False, SANS),
        ],
        [
            ("One hint pulls ", GREY, False, SANS),
            ("4 KB", GREEN, True, MONO),
            (" ahead — enough to hide a ", GREY, False, SANS),
            ("~1000-cycle", RED, True, MONO),
            (" I-cache miss.", GREY, False, SANS),
        ],
    ],
    size=16,
    gap=12,
)

# =========================================================================
# 5. GLOBAL BYTE GRID
# =========================================================================
s = slide()
header(s, "MODEL", "Global byte grid:  P(k) = 32640 + k x 4096", 7)
tf = txt(s, Inches(0.58), Inches(1.55), Inches(12.2), Inches(0.4))
para(
    tf,
    [
        (
            "Each grid step = one 4 KB hint (32 cache lines x 128 B). Anchors for every prefetch family.",
            DIM,
            False,
            SANS,
        )
    ],
    size=14,
    space_after=0,
    first=True,
)
# ---- CP block + repeating 4 KB grid cells ----
gy, gh = Inches(2.55), Inches(0.95)
cp = rect(s, Inches(0.8), gy, Inches(3.4), gh, fill=GREEN)
tfc = cp.text_frame
tfc.vertical_anchor = MSO_ANCHOR.MIDDLE
tfc.word_wrap = True
para(
    tfc,
    [("CP window\n[0, 32640 B)", INK, True, SANS)],
    size=13,
    align=PP_ALIGN.CENTER,
    space_after=0,
    first=True,
)
ncell = 7
cw = (Inches(12.5) - Inches(4.3)) / ncell
for k in range(ncell):
    cx = Inches(4.3) + cw * k
    cell = rect(s, cx, gy, cw, gh, fill=PANEL2, line=AMBER, line_w=1.0)
    tfx = cell.text_frame
    tfx.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(
        tfx,
        [("P%d" % k, AMBER, True, MONO)],
        size=13,
        align=PP_ALIGN.CENTER,
        space_after=1,
        first=True,
    )
    para(
        tfx, [("4 KB", DIM, False, MONO)], size=9, align=PP_ALIGN.CENTER, space_after=0
    )
# span bracket label over cells
tf = txt(s, Inches(4.3), gy - Inches(0.4), Inches(8.2), Inches(0.35))
para(
    tf,
    [("SW prefetch grid — one hint per P(k)", AMBER, True, SANS)],
    size=13,
    align=PP_ALIGN.CENTER,
    space_after=0,
    first=True,
)
# byte scale under key boundaries
for xx, lbl in [
    (Inches(4.3), "32640"),
    (Inches(4.3) + cw, "36736"),
    (Inches(4.3) + cw * 2, "40832"),
    (Inches(12.5), "65536"),
]:
    t = txt(s, xx - Inches(0.7), gy + gh + Inches(0.05), Inches(1.4), Inches(0.3))
    para(
        t,
        [(lbl, GREY, False, MONO)],
        size=11,
        align=PP_ALIGN.CENTER,
        space_after=0,
        first=True,
    )
# sequence line
tf = txt(s, Inches(0.8), Inches(4.35), Inches(11.9), Inches(0.4))
para(
    tf,
    [
        (
            "P(k) = 32640, 36736, 40832, 44928, 49024, 53120, 57216, …  (step 4096 B)",
            WHITE,
            True,
            MONO,
        )
    ],
    size=15,
    space_after=0,
    first=True,
)
# key rows
rows = [
    (
        "P(0) = 255 x 128 = 32640 B",
        "CP hardware max (.amdhsa_inst_pref_size saturates at 255)",
        GREEN,
    ),
    ("spacing = 4096 B", "32 cache lines x 128 B = one s_prefetch_inst hint", BLUE),
    ("I-cache = 65536 B", "static (<=64 KB) vs dynamic (>64 KB) split", AMBER),
]
y = Inches(4.95)
for a, b, col in rows:
    rect(s, Inches(0.58), y, Inches(0.14), Inches(0.5), fill=col)
    bx = card(s, Inches(0.85), y, Inches(11.9), Inches(0.5))
    tfr = bx.text_frame
    tfr.vertical_anchor = MSO_ANCHOR.MIDDLE
    tfr.margin_left = Pt(12)
    para(
        tfr,
        [(a, col, True, MONO), ("     " + b, GREY, False, SANS)],
        size=14,
        space_after=0,
        first=True,
    )
    y = y + Inches(0.6)

# =========================================================================
# 6. TAXONOMY
# =========================================================================
s = slide()
header(s, "OVERVIEW", "Prefetch taxonomy", 4)
card(s, Inches(0.55), Inches(1.65), Inches(3.75), Inches(5.15))
tf = txt(s, Inches(0.8), Inches(1.82), Inches(3.3), Inches(0.9))
para(tf, [("CP Prefetch", WHITE, True, SANS)], size=22, space_after=2, first=True)
para(tf, [("Hardware preloads the head", GREY, False, SANS)], size=13, space_after=0)
pill(
    s,
    Inches(0.8),
    Inches(3.0),
    Inches(3.25),
    "AccumulateInstructionSizePass",
    GREEN,
    size=10,
)
tf = txt(s, Inches(0.8), Inches(3.7), Inches(3.3), Inches(2.9))
para(
    tf,
    [("Covers ", GREY, False, SANS), ("[0, 32640 B)", GREEN, True, MONO)],
    size=15,
    space_after=8,
    first=True,
)
para(
    tf,
    [("Measures layout — inserts nothing.", GREY, False, SANS)],
    size=14,
    space_after=0,
    bullet="·",
)
card(s, Inches(4.55), Inches(1.65), Inches(8.25), Inches(5.15))
tf = txt(s, Inches(4.8), Inches(1.82), Inches(7.7), Inches(0.9))
para(tf, [("SW Prefetch", WHITE, True, SANS)], size=22, space_after=2, first=True)
para(
    tf,
    [("s_prefetch_inst hints beyond the CP boundary", GREY, False, SANS)],
    size=13,
    space_after=0,
)
card(s, Inches(4.8), Inches(2.9), Inches(3.85), Inches(3.7), fill=PANEL2)
tf = txt(s, Inches(5.0), Inches(3.02), Inches(3.45), Inches(0.5))
para(tf, [("Relative  (PC-rel)", BLUE, True, SANS)], size=17, space_after=0, first=True)
pill(
    s,
    Inches(5.0),
    Inches(3.6),
    Inches(3.45),
    "SwInstructionPrefetchRelStaticPass",
    BLUE,
    size=9,
)
pill(
    s,
    Inches(5.0),
    Inches(4.14),
    Inches(3.45),
    "SwInstructionPrefetchRelDynamicPass",
    BLUE,
    size=9,
)
tf = txt(s, Inches(5.0), Inches(4.78), Inches(3.45), Inches(1.7))
para(
    tf,
    [("The safe, universal default.", GREY, False, SANS)],
    size=14,
    space_after=6,
    first=True,
    bullet=None,
)
para(
    tf,
    [("All kernels · no SGPRs · small delay", DIM, False, SANS)],
    size=12,
    space_after=0,
    bullet=None,
)
card(s, Inches(8.9), Inches(2.9), Inches(3.65), Inches(3.7), fill=PANEL2)
tf = txt(s, Inches(9.1), Inches(3.02), Inches(3.25), Inches(0.5))
para(tf, [("Absolute  (getpc)", AMBER, True, SANS)], size=17, space_after=0, first=True)
pill(
    s,
    Inches(9.1),
    Inches(3.6),
    Inches(3.25),
    "SwInstructionPrefetchAbsStaticPass",
    AMBER,
    size=9,
)
pill(
    s,
    Inches(9.1),
    Inches(4.14),
    Inches(3.25),
    "SwInstructionPrefetchAbsDynamicPass",
    AMBER,
    size=9,
)
tf = txt(s, Inches(9.1), Inches(4.78), Inches(3.25), Inches(1.7))
para(
    tf,
    [("The low-latency, focused option.", GREY, False, SANS)],
    size=14,
    space_after=6,
    first=True,
    bullet=None,
)
para(
    tf,
    [("No delay · 3 SGPRs (hidden) · GEMM today", DIM, False, SANS)],
    size=12,
    space_after=0,
    bullet=None,
)

# =========================================================================
# 7. CP PREFETCH
# =========================================================================
s = slide()
header(s, "CP PREFETCH", "Command-processor prefetch", 5)
pill(
    s,
    Inches(0.58),
    Inches(1.7),
    Inches(4.3),
    "AccumulateInstructionSizePass",
    GREEN,
    h=0.5,
    size=13,
)
bullets(
    s,
    Inches(0.58),
    Inches(2.5),
    Inches(6.9),
    Inches(4.2),
    [
        [
            (
                "Preloads the kernel head into the I-cache before launch.",
                GREY,
                False,
                SANS,
            )
        ],
        [
            ("Prefetch is ", GREY, False, SANS),
            ("byte-based", WHITE, True, SANS),
            (
                " — SW hints only help if placed at the exact byte boundary.",
                GREY,
                False,
                SANS,
            ),
        ],
        [
            ("Driver prefetch is ", GREY, False, SANS),
            ("capped at 32640 B", GREEN, True, MONO),
            ("; oversize can error.", GREY, False, SANS),
        ],
        [
            ("Measures layout & builds the ", GREY, False, SANS),
            ("label -> byte-offset map", WHITE, True, SANS),
            (" (inserts nothing).", GREY, False, SANS),
        ],
    ],
    size=16,
    gap=14,
)
card(s, Inches(7.75), Inches(1.9), Inches(5.0), Inches(1.85), line=GREEN)
tf = txt(s, Inches(8.0), Inches(2.02), Inches(4.55), Inches(1.65))
para(
    tf,
    [("CP BOUNDARY  ", DIM, True, MONO), (".amdhsa_inst_pref_size", GREEN, True, MONO)],
    size=13,
    space_after=5,
    first=True,
)
para(
    tf,
    [
        ("[0, 32640 B)", GREEN, True, MONO),
        (" resident at launch — zero SW cost.", GREY, False, SANS),
    ],
    size=13,
    space_after=4,
)
para(
    tf,
    [
        ("unit 128 B · max 255 · ", GREY, False, SANS),
        ("255 x 128 = 32640", WHITE, True, MONO),
    ],
    size=13,
    space_after=4,
)
para(
    tf,
    [
        ("set > kernel .text -> ", GREY, False, SANS),
        ("driver may error", RED, True, SANS),
    ],
    size=13,
    space_after=0,
)
card(s, Inches(7.75), Inches(3.9), Inches(5.0), Inches(2.55))
tf = txt(s, Inches(8.0), Inches(4.02), Inches(4.5), Inches(2.35))
para(
    tf,
    [("BYTE-ACCURATE SIZING accounts for", DIM, True, MONO)],
    size=12,
    space_after=8,
    first=True,
)
for a, b in [
    ("base encoding bytes", ""),
    ("VALU ", "4 -> 8 B promotion"),
    ("literal ", "+4 / +8 B  (label / .set aware)"),
    (".align ", "+ label padding"),
]:
    para(
        tf,
        [
            (a, WHITE, True, SANS),
            (b, AMBER if b else WHITE, b != "", MONO if b else SANS),
        ],
        size=14,
        space_after=7,
        bullet="•",
    )
status_line(
    s,
    Inches(0.58),
    Inches(6.35),
    Inches(7.0),
    status="in codegen — shipping",
    challenge="byte-exact sizing across ISA groups, VALU 4->8, literals, .align",
)


# =========================================================================
# family helper (slides 8, 10)
# =========================================================================
def family_slide(
    kicker, title, page, accent, pros, cons, passes, status=None, challenge=None
):
    s = slide()
    header(s, kicker, title, page)

    def _items(tf, items, sym):
        for it in items:
            main, sub = it if isinstance(it, tuple) else (it, None)
            para(
                tf,
                [(main, WHITE, True, SANS)],
                size=15,
                space_after=(1 if sub else 8),
                bullet=sym,
            )
            if sub:
                para(tf, [(sub, GREY, False, SANS)], size=13, space_after=8, level=1)

    card(s, Inches(0.58), Inches(1.7), Inches(5.9), Inches(2.15), line=GREEN)
    tf = txt(s, Inches(0.8), Inches(1.84), Inches(5.5), Inches(1.9))
    para(tf, [("PROS", GREEN, True, MONO)], size=13, space_after=8, first=True)
    _items(tf, pros, "+")
    card(s, Inches(0.58), Inches(4.0), Inches(5.9), Inches(2.6), line=RED)
    tf = txt(s, Inches(0.8), Inches(4.14), Inches(5.5), Inches(2.35))
    para(tf, [("CONS", RED, True, MONO)], size=13, space_after=8, first=True)
    _items(tf, cons, "-")
    y = Inches(1.7)
    for name, sub in passes:
        pill(s, Inches(6.8), y, Inches(5.95), name, accent, size=13, h=0.5)
        tf = txt(s, Inches(7.0), y + Inches(0.58), Inches(5.7), Inches(1.1))
        for i, ln in enumerate(sub):
            para(
                tf,
                [(ln, GREY, False, SANS)],
                size=14,
                space_after=3,
                bullet="·",
                first=(i == 0),
            )
        y = y + Inches(0.66 + 0.32 * len(sub) + 0.26)
    status_line(s, Inches(6.8), Inches(5.55), Inches(5.95), status, challenge)
    return s


# =========================================================================
# 8. RELATIVE OVERVIEW
# =========================================================================
family_slide(
    "RELATIVE",
    "Relative prefetch (PC-relative)",
    8,
    BLUE,
    pros=[
        ("No extra SGPRs", None),
        (
            "Full coverage",
            "Pure GEMM \u00b7 Fused GEMM \u00b7 StreamK \u00b7 Sparse GEMM",
        ),
    ],
    cons=[("Data-transfer delay", "PC-relative address resolved at prefetch time")],
    passes=[
        (
            "SwInstructionPrefetchRelStaticPass",
            [
                "Insertion by static kernel-layout byte offset",
                "Kept for stinkytofu-opt / unit tests",
            ],
        ),
        (
            "SwInstructionPrefetchRelDynamicPass",
            [
                "ControlFlowGraph (CFG)-gated, per-basic-block insertion",
                "Wider coverage -> better performance",
                "Default pass wired on gfx1250",
            ],
        ),
    ],
    status="RelStatic shipping in codegen; RelDynamic done, in review",
    challenge="CFG-gated per-path byte tracking; skip loop bodies",
)

# =========================================================================
# 9. RELATIVE: LAYOUT vs CFG
# =========================================================================
s = slide()
header(s, "RELATIVE - CONCEPT", "Relative: grid-based insertion", 9)
# ---- grid-insertion strip (left) ----
sx, sy, cw2, ch2 = Inches(0.75), Inches(2.9), Inches(0.66), Inches(0.62)
ncell = 8
# BB labels
for lbl, cx0, span in [("BB1", sx, 4), ("BB2", sx + cw2 * 4, 4)]:
    t = txt(s, cx0, sy - Inches(0.42), cw2 * span, Inches(0.3))
    para(
        t,
        [(lbl, DIM, True, MONO)],
        size=12,
        align=PP_ALIGN.CENTER,
        space_after=0,
        first=True,
    )
for i in range(ncell):
    cx = sx + cw2 * i
    tgt = i in (3, 6)
    cell = rect(s, cx, sy, cw2, ch2, fill=PANEL2 if tgt else DARK, line=LINEC)
    tfx = cell.text_frame
    tfx.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(
        tfx,
        [("I%d" % i, GREY, False, MONO)],
        size=11,
        align=PP_ALIGN.CENTER,
        space_after=0,
        first=True,
    )
    if tgt:
        rect(s, cx, sy, Inches(0.07), ch2, fill=BLUE)  # "inserted before" marker
# divider between BB1 and BB2
rect(s, sx + cw2 * 4, sy - Inches(0.12), Inches(0.02), ch2 + Inches(0.24), fill=DIM)
# P(k) ticks (red) at target left edges + labels above
for tcell, lbl in [(3, "P(k)"), (6, "P(k+1)")]:
    cx = sx + cw2 * tcell
    rect(
        s,
        cx - Inches(0.01),
        sy - Inches(0.55),
        Inches(0.03),
        ch2 + Inches(0.75),
        fill=RED,
    )
    t = txt(s, cx - Inches(0.55), sy - Inches(0.95), Inches(1.1), Inches(0.32))
    para(
        t,
        [(lbl, RED, True, MONO)],
        size=12,
        align=PP_ALIGN.CENTER,
        space_after=0,
        first=True,
    )
# 4096 B span label
t = txt(s, sx + cw2 * 3, sy + ch2 + Inches(0.08), cw2 * 3, Inches(0.3))
para(
    t,
    [("|<-- 4096 B -->|", DIM, False, MONO)],
    size=12,
    align=PP_ALIGN.CENTER,
    space_after=0,
    first=True,
)
# SWP tags below targets
for tcell in (3, 6):
    cx = sx + cw2 * tcell
    t = txt(s, cx - Inches(0.35), sy + ch2 + Inches(0.5), Inches(1.3), Inches(0.3))
    para(
        t,
        [("SWP", BLUE, True, MONO)],
        size=12,
        align=PP_ALIGN.LEFT,
        space_after=0,
        first=True,
    )
card(s, Inches(0.58), Inches(4.75), Inches(6.1), Inches(0.7), fill=PANEL2)
tf = txt(s, Inches(0.78), Inches(4.75), Inches(5.7), Inches(0.7), MSO_ANCHOR.MIDDLE)
para(
    tf,
    [
        ("Rule:  ", DIM, True, MONO),
        ("globalPcBefore < P(k) <= globalPcAfter", WHITE, True, MONO),
    ],
    size=13,
    space_after=0,
    first=True,
)
tf = txt(s, Inches(0.58), Inches(5.65), Inches(6.1), Inches(1.2))
para(
    tf,
    [
        (
            "Insert the hint before the instruction whose byte span contains P(k).",
            GREY,
            False,
            SANS,
        )
    ],
    size=14,
    space_after=0,
    first=True,
)
# ---- bullets (right) ----
bullets(
    s,
    Inches(7.0),
    Inches(1.85),
    Inches(5.8),
    Inches(3.0),
    [
        [
            ("Static", AMBER, True, SANS),
            (" walks ", GREY, False, SANS),
            ("layout order", WHITE, True, SANS),
            (" — blind to control flow.", GREY, False, SANS),
        ],
        [
            ("Dynamic", GREEN, True, SANS),
            (" tracks ", GREY, False, SANS),
            ("execution progress", WHITE, True, SANS),
            (" per path; ", GREY, False, SANS),
            ("both branch arms", WHITE, True, SANS),
            (" get hints.", GREY, False, SANS),
        ],
        [
            ("Loops: ", GREY, False, SANS),
            ("back-edge ignored", WHITE, True, SANS),
            (" — body cached after iteration 1.", GREY, False, SANS),
        ],
    ],
    size=15,
    gap=12,
)


# ---- STATIC vs DYNAMIC branch-coverage mini-CFG (bottom-right) ----
def _mini_cfg(x, label, lcol, arm2_covered):
    t = txt(s, x, Inches(4.75), Inches(2.7), Inches(0.28))
    para(
        t,
        [(label, lcol, True, MONO)],
        size=11,
        align=PP_ALIGN.CENTER,
        space_after=0,
        first=True,
    )
    e = rect(
        s,
        x + Inches(0.7),
        Inches(5.1),
        Inches(1.3),
        Inches(0.34),
        fill=DARK,
        line=LINEC,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE,
    )
    te = e.text_frame
    te.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(
        te,
        [("BB0", GREY, True, MONO)],
        size=10,
        align=PP_ALIGN.CENTER,
        space_after=0,
        first=True,
    )
    rect(s, x + Inches(0.65), Inches(5.58), Inches(1.4), Inches(0.02), fill=DIM)
    rect(s, x + Inches(0.65), Inches(5.58), Inches(0.02), Inches(0.14), fill=DIM)
    rect(s, x + Inches(2.03), Inches(5.58), Inches(0.02), Inches(0.14), fill=DIM)
    a1 = rect(
        s,
        x,
        Inches(5.72),
        Inches(1.3),
        Inches(0.38),
        fill=PANEL2,
        line=GREEN,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE,
    )
    ta = a1.text_frame
    ta.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(
        ta,
        [("BB1 · SWP", GREEN, True, MONO)],
        size=10,
        align=PP_ALIGN.CENTER,
        space_after=0,
        first=True,
    )
    c2 = GREEN if arm2_covered else RED
    a2 = rect(
        s,
        x + Inches(1.4),
        Inches(5.72),
        Inches(1.3),
        Inches(0.38),
        fill=PANEL2,
        line=c2,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE,
    )
    tb = a2.text_frame
    tb.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(
        tb,
        [("BB2 · SWP" if arm2_covered else "BB2 · miss", c2, True, MONO)],
        size=10,
        align=PP_ALIGN.CENTER,
        space_after=0,
        first=True,
    )


def _static_tape(x):
    t = txt(s, x, Inches(4.75), Inches(2.7), Inches(0.28))
    para(
        t,
        [("STATIC — linear layout", AMBER, True, MONO)],
        size=11,
        align=PP_ALIGN.CENTER,
        space_after=0,
        first=True,
    )
    ta = txt(s, x + Inches(0.87), Inches(5.08), Inches(0.9), Inches(0.24))
    para(
        ta,
        [("P(k) hit", GREEN, True, MONO)],
        size=9,
        align=PP_ALIGN.CENTER,
        space_after=0,
        first=True,
    )
    tb = txt(s, x + Inches(1.72), Inches(5.08), Inches(0.98), Inches(0.24))
    para(
        tb,
        [("off-grid", RED, True, MONO)],
        size=9,
        align=PP_ALIGN.CENTER,
        space_after=0,
        first=True,
    )
    for i, (lbl, col) in enumerate([("BB0", LINEC), ("BB1", GREEN), ("BB2", RED)]):
        bx = rect(
            s,
            x + Inches(0.02 + 0.89 * i),
            Inches(5.35),
            Inches(0.85),
            Inches(0.4),
            fill=DARK if i == 0 else PANEL2,
            line=col,
            line_w=1.25,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        )
        tbx = bx.text_frame
        tbx.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(
            tbx,
            [(lbl, GREY if i == 0 else col, True, MONO)],
            size=10,
            align=PP_ALIGN.CENTER,
            space_after=0,
            first=True,
        )


_static_tape(Inches(7.0))
_mini_cfg(Inches(10.0), "DYNAMIC — per-BB CFG", GREEN, True)
t = txt(s, Inches(7.0), Inches(6.22), Inches(5.75), Inches(0.75))
para(
    t,
    [
        ("Static reads code as a flat tape — ", GREY, False, SANS),
        ("BB2 is a branch target off the grid", WHITE, True, SANS),
        (", so it is missed. Dynamic restarts per BB -> ", GREY, False, SANS),
        ("both BB1 & BB2 hinted", GREEN, True, SANS),
        (".", GREY, False, SANS),
    ],
    size=12,
    space_after=0,
    first=True,
)

# =========================================================================
# 10. ABSOLUTE OVERVIEW
# =========================================================================
family_slide(
    "ABSOLUTE",
    "Absolute prefetch (getpc base)",
    10,
    AMBER,
    pros=[
        (
            "No data-transfer delay",
            "Absolute base address is ready \u2014 no per-hint compute",
        )
    ],
    cons=[
        ("3 extra SGPRs", "Can hide with register allocation \u2014 Done"),
        ("Narrow coverage", "No support on Fused GEMM / StreamK currently"),
    ],
    passes=[
        (
            "SwInstructionPrefetchAbsStaticPass",
            [
                "32640 B < total <= 64 KB (fits I-cache)",
                "One entry burst: single label + koffsets",
            ],
        ),
        (
            "SwInstructionPrefetchAbsDynamicPass",
            [
                "total > 64 KB (kernel streams)",
                "ControlFlowGraph (CFG)-target ladder into hot blocks",
                "6 hints per arm (~24 KB prefetched ahead)",
            ],
        ),
    ],
    status="Implemented; done — in review (extensive experiments ongoing)",
    challenge="3 SGPRs hidden via reg-alloc; coverage limited to GEMM today",
)

# =========================================================================
# 11. ABS STATIC BURST
# =========================================================================
s = slide()
header(s, "ABSOLUTE - 1", "AbsStatic: single-label + koffset burst", 12)
pill(
    s,
    Inches(0.58),
    Inches(1.7),
    Inches(5.2),
    "SwInstructionPrefetchAbsStaticPass",
    AMBER,
    h=0.5,
    size=13,
)
bullets(
    s,
    Inches(0.58),
    Inches(2.45),
    Inches(6.5),
    Inches(4.4),
    [
        [
            (
                "One entry burst before any branch; whole kernel fits the I-cache.",
                GREY,
                False,
                SANS,
            )
        ],
        [
            ("One label at ", GREY, False, SANS),
            ("P(0)", AMBER, True, MONO),
            ("; koffsets 0, 4096, 8192 ... reuse the same base.", GREY, False, SANS),
        ],
        [
            ("getpc + bare-label reloc", WHITE, True, SANS),
            (" (+4 corrects for getpc).", GREY, False, SANS),
        ],
        [
            ("N solved by fixed-point; capped to I-cache: ", GREY, False, SANS),
            ("N <= (65536-32640)/4096 = 8", GREEN, True, MONO),
            (".", GREY, False, SANS),
        ],
    ],
    size=15,
    gap=12,
)
codebox(
    s,
    Inches(7.35),
    Inches(2.45),
    Inches(5.4),
    Inches(4.0),
    title="EMITTED BURST",
    lines=[
        ("label_Do_SW_PrefetchAbs_entry:", AMBER, True, MONO),
        ("s_getpc_b64 s[base:base+1]", WHITE, False, MONO),
        ("s_add_i32   s[base+2], L0, 4", GREEN, False, MONO),
        ("s_add_u32   s[base],   s[base], s[base+2]", WHITE, False, MONO),
        ("s_addc_u32  s[base+1], s[base+1], 0", WHITE, False, MONO),
        ("s_prefetch_inst s[base:base+1],", WHITE, False, MONO),
        ("            k*4096, null, 0x1f   x N", WHITE, False, MONO),
        ("...", DIM, False, MONO),
        ("L0 = label_SW_PrefetchAbs_0 @ P(0)", DIM, False, MONO),
        ("base:base+1 = addr;  base+2 = scratch", DIM, False, MONO),
    ],
    size=11,
)
# entry-burst BB strip (left): one burst at entry, single target label L0 @ P(0)
_e = rect(
    s,
    Inches(0.58),
    Inches(4.8),
    Inches(2.25),
    Inches(0.72),
    fill=PANEL2,
    line=AMBER,
    line_w=1.25,
    shape=MSO_SHAPE.ROUNDED_RECTANGLE,
)
_t = _e.text_frame
_t.vertical_anchor = MSO_ANCHOR.MIDDLE
_t.word_wrap = True
para(
    _t,
    [("entry BB\n(getpc + burst)", AMBER, True, MONO)],
    size=11,
    align=PP_ALIGN.CENTER,
    space_after=0,
    first=True,
)
rect(s, Inches(2.83), Inches(5.13), Inches(0.32), Inches(0.03), fill=DIM)
_b = rect(
    s,
    Inches(3.15),
    Inches(4.8),
    Inches(3.9),
    Inches(0.72),
    fill=DARK,
    line=LINEC,
    shape=MSO_SHAPE.ROUNDED_RECTANGLE,
)
_t = _b.text_frame
_t.vertical_anchor = MSO_ANCHOR.MIDDLE
_t.word_wrap = True
para(
    _t,
    [("kernel body", WHITE, True, MONO), ("   L0 @ P(0)", GREEN, True, MONO)],
    size=11,
    align=PP_ALIGN.CENTER,
    space_after=0,
    first=True,
)
_t = txt(s, Inches(0.58), Inches(5.6), Inches(6.5), Inches(0.4))
para(
    _t,
    [
        (
            "One burst at entry; koffsets 0, 4K, 8K … reach the body from L0.",
            GREY,
            False,
            SANS,
        )
    ],
    size=12,
    space_after=0,
    first=True,
)
status_line(
    s,
    Inches(0.58),
    Inches(6.1),
    Inches(6.6),
    status="done — in review",
    challenge="getpc + label math; fixed-point N vs post-insert shift",
)

# =========================================================================
# 12. ABS DYNAMIC
# =========================================================================
s = slide()
header(s, "ABSOLUTE - 2", "AbsDynamic: CFG-target ladder (> 64 KB)", 13)
pill(
    s,
    Inches(0.58),
    Inches(1.7),
    Inches(5.4),
    "SwInstructionPrefetchAbsDynamicPass",
    AMBER,
    h=0.5,
    size=13,
)
bullets(
    s,
    Inches(0.58),
    Inches(2.45),
    Inches(7.0),
    Inches(3.6),
    [
        [
            ("The kernel is ", GREY, False, SANS),
            ("too big for the I-cache", WHITE, True, SANS),
            (
                " — it streams, so one entry burst can't keep it all hot.",
                GREY,
                False,
                SANS,
            ),
        ],
        [
            ("Instead, prefetch only the block that runs ", GREY, False, SANS),
            ("next", WHITE, True, SANS),
            (": the hot ", GREY, False, SANS),
            ("global-write (output) stage", WHITE, True, SANS),
            (".", GREY, False, SANS),
        ],
        [
            ("Which output block runs depends on ", GREY, False, SANS),
            ("runtime args (GSU, Beta)", WHITE, True, SANS),
            (".", GREY, False, SANS),
        ],
        [
            ("Insert the hints once in the ", GREY, False, SANS),
            ("preloop", AMBER, True, SANS),
            ("; each target gets ", GREY, False, SANS),
            ("6 hints (~24 KB)", GREEN, True, MONO),
            (" ahead.", GREY, False, SANS),
        ],
    ],
    size=15,
    gap=13,
)
status_line(
    s,
    Inches(0.58),
    Inches(6.0),
    Inches(7.0),
    status="landed & enabled, pending gfx1250 numeric validation",
    challenge="choose the right output block from (GSU, Beta) at run time",
)
# BB target-selection graph (right): MultiGemmEnd -> runtime (GSU,Beta) -> hot GW block
card(s, Inches(7.75), Inches(2.35), Inches(5.0), Inches(3.55))
t = txt(s, Inches(7.95), Inches(2.44), Inches(4.6), Inches(0.28))
para(
    t,
    [("TARGET SELECTION (runtime)", DIM, True, MONO)],
    size=11,
    space_after=0,
    first=True,
)


def _node(x, y, w, h, text, col, fill=PANEL2):
    b = rect(
        s,
        x,
        y,
        w,
        h,
        fill=fill,
        line=col,
        line_w=1.25,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE,
    )
    tb = b.text_frame
    tb.vertical_anchor = MSO_ANCHOR.MIDDLE
    tb.word_wrap = True
    para(
        tb,
        [(text, col, True, MONO)],
        size=12,
        align=PP_ALIGN.CENTER,
        space_after=0,
        first=True,
    )


_node(Inches(8.45), Inches(2.85), Inches(3.6), Inches(0.5), "preloop (once)", AMBER)
rect(s, Inches(10.23), Inches(3.35), Inches(0.03), Inches(0.28), fill=DIM)
_node(
    Inches(8.45),
    Inches(3.63),
    Inches(3.6),
    Inches(0.5),
    "runtime: (GSU, Beta) ?",
    WHITE,
    DARK,
)
rect(s, Inches(9.15), Inches(4.13), Inches(2.2), Inches(0.03), fill=DIM)  # split bar
rect(s, Inches(9.15), Inches(4.13), Inches(0.03), Inches(0.28), fill=DIM)
rect(s, Inches(11.32), Inches(4.13), Inches(0.03), Inches(0.28), fill=DIM)
_node(Inches(8.05), Inches(4.4), Inches(2.1), Inches(0.55), "GW_B0_GSU1", GREEN)
_node(Inches(10.35), Inches(4.4), Inches(2.1), Inches(0.55), "GW_B1_GSU1", GREEN)
t = txt(s, Inches(7.95), Inches(5.15), Inches(4.65), Inches(0.65))
para(
    t,
    [("Prefetch the output block the kernel will run next.", GREY, False, SANS)],
    size=12,
    space_after=0,
    first=True,
)

# =========================================================================
# 13. SIZE-BASED SELECTION + PIPELINE
# =========================================================================
s = slide()
header(s, "POLICY", "Size-based selection & pipeline", 11)
tf = txt(s, Inches(0.8), Inches(1.55), Inches(11.9), Inches(0.4))
para(
    tf,
    [
        ("Family first: ", DIM, False, SANS),
        ("Absolute > PC-rel", AMBER, True, SANS),
        (
            " (mutually exclusive, default PC-rel). Absolute then picks a tier by total .text size:",
            DIM,
            False,
            SANS,
        ),
    ],
    size=13,
    space_after=0,
    first=True,
)
axis_y = Inches(3.35)
rect(s, Inches(0.8), axis_y, Inches(11.7), Inches(0.05), fill=DIM)
segs = [
    (
        Inches(0.8),
        Inches(3.4),
        GREEN,
        "0 - 32640 B",
        "CP Prefetch",
        "CP preload covers the head",
    ),
    (
        Inches(4.2),
        Inches(4.6),
        AMBER,
        "32640 B - 64 KB",
        "AbsStaticPass",
        "Fits I-cache . single burst",
    ),
    (
        Inches(8.8),
        Inches(3.7),
        RED,
        ">  64 KB",
        "AbsDynamicPass",
        "Streams . CFG-target ladder",
    ),
]
for x, w, color, rng, name, desc in segs:
    box = rect(
        s,
        x,
        Inches(2.45),
        w - Inches(0.15),
        Inches(0.85),
        fill=color,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE,
    )
    tfb = box.text_frame
    tfb.word_wrap = True
    tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(
        tfb,
        [(rng, INK, True, MONO)],
        size=16,
        align=PP_ALIGN.CENTER,
        space_after=0,
        first=True,
    )
    rect(s, x, axis_y, Inches(0.03), Inches(0.2), fill=color)
    card(s, x, Inches(3.95), w - Inches(0.15), Inches(1.4), line=color, line_w=1.25)
    tf2 = txt(s, x + Inches(0.2), Inches(4.08), w - Inches(0.55), Inches(1.2))
    ncol = AMBER if name.startswith("Abs") else color
    para(tf2, [(name, ncol, True, MONO)], size=15, space_after=6, first=True)
    para(tf2, [(desc, GREY, False, SANS)], size=13, space_after=0, line_spacing=1.05)
card(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.15), fill=PANEL2)
tf = txt(s, Inches(1.0), Inches(5.72), Inches(11.4), Inches(0.95), MSO_ANCHOR.MIDDLE)
para(
    tf,
    [
        ("PIPELINE   ", DIM, True, MONO),
        (
            "SetMatrixReuse -> Abs(Dynamic+Static) or RelDynamic -> AccumInstrSize",
            WHITE,
            False,
            MONO,
        ),
    ],
    size=14,
    space_after=2,
    first=True,
)
para(
    tf,
    [
        (
            "Both absolute passes are added; each no-ops outside its size regime. Default = RelDynamic.",
            GREY,
            False,
            SANS,
        )
    ],
    size=13,
    space_after=0,
)

# =========================================================================
# 14. COMPARISON TABLE
# =========================================================================
s = slide()
header(s, "SUMMARY", "Relative vs Absolute", 14)
rows = [
    ("", "Relative (PC-rel)", "Absolute (getpc)"),
    ("Instruction", "s_prefetch_inst_pc_rel", "s_prefetch_inst"),
    ("Address", "PC-relative", "getpc + bare-label reloc"),
    ("Extra SGPRs", "None", "3 (pair + scratch, hidden)"),
    ("Transfer delay", "Yes", "No"),
    ("Coverage", "GEMM . Fused . StreamK . Sparse", "GEMM only (today)"),
    ("Passes", "RelStatic / RelDynamic", "AbsStatic / AbsDynamic"),
    ("Best when", "Broad coverage & correctness", "Latency-critical, SGPRs free"),
]
x0, y0 = Inches(0.58), Inches(1.7)
cw = [Inches(3.0), Inches(4.65), Inches(4.65)]
rh = Inches(0.6)
colors = [None, BLUE, AMBER]
for ri, row in enumerate(rows):
    cx = x0
    for ci, cell in enumerate(row):
        head = ri == 0
        fill = PANEL2 if head else (PANEL if ri % 2 else DARK)
        cs = rect(
            s, cx, y0 + rh * ri, cw[ci], rh, fill=fill, line=CELL_LINE, line_w=0.75
        )
        tfr = cs.text_frame
        tfr.vertical_anchor = MSO_ANCHOR.MIDDLE
        tfr.margin_left = Pt(8)
        if head:
            para(
                tfr,
                [(cell, colors[ci] or WHITE, True, MONO)],
                size=16,
                space_after=0,
                first=True,
            )
        elif ci == 0:
            para(tfr, [(cell, GREY, True, SANS)], size=14, space_after=0, first=True)
        else:
            para(tfr, [(cell, WHITE, False, SANS)], size=14, space_after=0, first=True)
        cx = cx + cw[ci]

# =========================================================================
# 15. CACHE-MISS LATENCY ACROSS PASSES
# =========================================================================
s = slide()
header(s, "EVIDENCE", "Cache-miss latency across the passes", 15)
tf = txt(s, Inches(0.58), Inches(1.45), Inches(12.2), Inches(0.4))
para(
    tf,
    [
        ("Instruction-request miss latency (SQC), ", DIM, False, SANS),
        ("from AM reports", WHITE, True, SANS),
        (
            " — same kernel, each variant. Lower & flatter post-CP plateau = fewer stalls.",
            DIM,
            False,
            SANS,
        ),
    ],
    size=13,
    space_after=0,
    first=True,
)
gw = Inches(3.05)
cols = [Inches(0.65), Inches(4.95), Inches(9.25)]
cells = [
    (0, 0, "evidence/miss_no_prefetch.png", "No prefetch", RED),
    (1, 0, "evidence/miss_original.png", "CP prefetch only", GREY),
    (2, 0, "evidence/miss_rel_static.png", "Relative Static", BLUE),
    (0, 1, "evidence/miss_rel_dyn_perbb.png", "Relative Dynamic", BLUE),
    (1, 1, "evidence/miss_abs_static.png", "Absolute Static", AMBER),
    (2, 1, "evidence/miss_abs_dynamic.png", "Absolute Dynamic (measured)", GREEN),
]
rowy = [Inches(1.95), Inches(4.3)]
lbly = [Inches(1.72), Inches(4.07)]
for c, r, path, label, col in cells:
    t = txt(s, cols[c], lbly[r], gw, Inches(0.26))
    para(t, [(label, col, True, SANS)], size=12, space_after=0, first=True)
    picture(s, path, cols[c], rowy[r], gw)
# takeaway (full-width line) + config
t = txt(s, Inches(0.58), Inches(6.3), Inches(12.2), Inches(0.55))
para(
    t,
    [
        (
            "Prefetch flattens the post-CP miss plateau (~500 cycles -> near-flat).  ",
            GREY,
            False,
            SANS,
        ),
        ("Absolute Dynamic (measured) is flattest", GREEN, True, SANS),
        (" — matches ", GREY, False, SANS),
        ("+59% on B5-1", GREEN, True, MONO),
        ("; CP-only already removes the head misses.", GREY, False, SANS),
    ],
    size=12,
    space_after=0,
    first=True,
)
tf = txt(s, Inches(0.58), Inches(6.9), Inches(12.2), Inches(0.28))
para(tf, [(KCFG, DIM, False, SANS)], size=10, space_after=0, first=True)

# =========================================================================
# 16. RESULTS
# =========================================================================
s = slide()
header(s, "RESULTS", "Measured performance (B5-1)", 16)
tf = txt(s, Inches(0.58), Inches(1.5), Inches(12.2), Inches(0.35))
para(
    tf,
    [
        ("Two scenarios: ", DIM, False, SANS),
        ("cache-resident", WHITE, True, SANS),
        (" (no miss — nothing to hide) vs ", DIM, False, SANS),
        ("Rotate-Object", WHITE, True, SANS),
        (
            " (real I-cache pressure). Prefetch only matters under pressure.",
            DIM,
            False,
            SANS,
        ),
    ],
    size=13,
    space_after=0,
    first=True,
)
# headline (original big-number style)
rect(
    s,
    Inches(0.58),
    Inches(1.9),
    Inches(4.6),
    Inches(1.0),
    fill=PANEL2,
    line=GREEN,
    line_w=1.5,
    shape=MSO_SHAPE.ROUNDED_RECTANGLE,
)
tfb = txt(s, Inches(0.58), Inches(1.93), Inches(4.6), Inches(0.94), MSO_ANCHOR.MIDDLE)
para(
    tfb,
    [("+59%", GREEN, True, SANS)],
    size=34,
    align=PP_ALIGN.CENTER,
    space_after=0,
    first=True,
)
para(
    tfb,
    [("CP + Absolute Dynamic (best)", GREY, False, SANS)],
    size=12,
    align=PP_ALIGN.CENTER,
    space_after=0,
)
tf = txt(s, Inches(5.5), Inches(1.93), Inches(7.2), Inches(0.94), MSO_ANCHOR.MIDDLE)
para(
    tf,
    [("91.96 us -> 57.83 us", WHITE, True, MONO)],
    size=20,
    space_after=2,
    first=True,
)
para(
    tf,
    [
        ("186.8 -> 294.9 TFLOP/s", GREEN, True, MONO),
        ("   (under I-cache pressure)", DIM, False, SANS),
    ],
    size=14,
    space_after=0,
)
# full table: 6 rows x (variant + 2 scenarios + uplift)
trows = [
    (
        "Variant",
        "us\n(no miss)",
        "TFLOP/s\n(no miss)",
        "us\n(pressure)",
        "TFLOP/s\n(pressure)",
        "Uplift\nvs no prefetch",
    ),
    ("Original (no prefetch)", "57.56", "298.5", "91.96", "186.8", "-"),
    ("CP prefetch only", "53.73", "319.7", "80.69", "212.9", "+14%"),
    ("CP + Relative Static", "53.75", "319.6", "59.91", "286.7", "+53%"),
    ("CP + Relative Dynamic  (default)", "53.76", "319.5", "58.18", "295.3", "+58%"),
    ("CP + Absolute Static", "53.34", "322.1", "80.49", "213.4", "+14%"),
    ("CP + Absolute Dynamic", "53.82", "319.2", "57.83", "294.9", "+59%  best"),
]
x0, y0 = Inches(0.58), Inches(3.0)
cw = [Inches(3.2), Inches(1.55), Inches(1.75), Inches(1.55), Inches(1.75), Inches(1.85)]
hh, rh = Inches(0.6), Inches(0.46)
best_ri = len(trows) - 1
for ri, row in enumerate(trows):
    head = ri == 0
    yy = y0 if ri == 0 else y0 + hh + rh * (ri - 1)
    rowh = hh if head else rh
    best = ri == best_ri
    default = ri == 4
    cx = x0
    for ci, cell in enumerate(row):
        fill = PANEL2 if head else (BEST_FILL if best else (PANEL if ri % 2 else DARK))
        cs = rect(s, cx, yy, cw[ci], rowh, fill=fill, line=CELL_LINE, line_w=0.75)
        tfr = cs.text_frame
        tfr.vertical_anchor = MSO_ANCHOR.MIDDLE
        tfr.margin_left = Pt(8)
        tfr.word_wrap = True
        if head:
            for li, line in enumerate(cell.split("\n")):
                para(
                    tfr,
                    [(line, DIM, True, SANS)],
                    size=11,
                    space_after=0,
                    first=(li == 0),
                    line_spacing=1.0,
                )
        elif ci == 0:
            vc = GREEN if best else (BLUE if default else GREY)
            para(tfr, [(cell, vc, True, SANS)], size=12, space_after=0, first=True)
        elif ci == 5:
            up = "+" in cell
            para(
                tfr,
                [
                    (
                        cell,
                        GREEN if (best or (up and cell not in ("+14%",))) else GREY,
                        True,
                        MONO,
                    )
                ],
                size=12,
                space_after=0,
                first=True,
            )
        else:
            para(tfr, [(cell, WHITE, False, MONO)], size=12, space_after=0, first=True)
        cx = cx + cw[ci]
tf = txt(s, Inches(0.58), Inches(6.45), Inches(12.2), Inches(0.32))
para(
    tf,
    [
        (
            "No-miss case: all variants within ~1% (cache-resident — nothing to hide). ",
            DIM,
            False,
            SANS,
        ),
        (
            "Relative Dynamic ships by default; Absolute Dynamic is the latency ceiling.",
            GREY,
            True,
            SANS,
        ),
    ],
    size=12,
    space_after=0,
    first=True,
)
tf = txt(s, Inches(0.58), Inches(6.82), Inches(12.2), Inches(0.3))
para(tf, [(KCFG, DIM, False, SANS)], size=10, space_after=0, first=True)

# --- Reorder slides to the presentation order (Overview5): ------------------
# creation order (0-based):
#   0 Title, 1 Agenda, 2 Motivation, 3 ISA, 4 ByteGrid, 5 Taxonomy, 6 CP,
#   7 RelOverview, 8 RelGrid, 9 AbsOverview, 10 AbsStatic, 11 AbsDynamic,
#   12 Policy, 13 Summary, 14 Evidence, 15 Results
# desired physical order:
#   Title, Agenda, Motivation, Taxonomy, CP, ByteGrid, ISA, RelOverview,
#   RelGrid, AbsOverview, Policy, AbsStatic, AbsDynamic, Summary, Evidence, Results
_order = [0, 1, 2, 5, 6, 3, 4, 7, 8, 9, 12, 10, 11, 13, 14, 15]
_lst = prs.slides._sldIdLst
_ids = list(_lst)
for _sid in _ids:
    _lst.remove(_sid)
for _i in _order:
    _lst.append(_ids[_i])

out = OUT
prs.save(out)
print("wrote", out, "with", len(prs.slides._sldIdLst), "slides")
