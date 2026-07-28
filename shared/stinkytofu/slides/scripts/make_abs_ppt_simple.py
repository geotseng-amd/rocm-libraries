#!/usr/bin/env python3
"""Single, simple, easy-to-read slide for abs dynamic prefetch coverage."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

BG = RGBColor(0x16, 0x16, 0x16)
FG = RGBColor(0xF0, 0xF0, 0xF0)
SUB = RGBColor(0xA8, 0xA8, 0xA8)
ACC = RGBColor(0x59, 0x9C, 0xE7)
GREEN = RGBColor(0x46, 0xB0, 0x72)
YELLOW = RGBColor(0xE2, 0xB5, 0x55)
RED = RGBColor(0xE0, 0x5A, 0x6E)
PANEL = RGBColor(0x20, 0x20, 0x20)
PANEL2 = RGBColor(0x24, 0x2C, 0x38)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid()
s.background.fill.fore_color.rgb = BG
SW = prs.slide_width


def run(p, t, sz, c=FG, b=False):
    r = p.add_run()
    r.text = t
    f = r.font
    f.size = Pt(sz)
    f.bold = b
    f.color.rgb = c
    f.name = "Calibri"
    return r


def box(l, t, w, h, fill=None):
    if fill is not None:
        sp = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
        sp.line.fill.background()
        tf = sp.text_frame
    else:
        tf = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    return tf


def dot(l, t, c, d=0.16):
    sp = s.shapes.add_shape(9, Inches(l), Inches(t), Inches(d), Inches(d))
    sp.fill.solid()
    sp.fill.fore_color.rgb = c
    sp.line.fill.background()


# Title
bar = s.shapes.add_shape(1, 0, 0, SW, Inches(0.10))
bar.fill.solid()
bar.fill.fore_color.rgb = ACC
bar.line.fill.background()
tf = box(0.45, 0.22, 12.5, 1.0)
run(
    tf.paragraphs[0],
    "Does abs Instruction Prefetch cover the GEMM epilogue? (gfx1250)",
    24,
    FG,
    b=True,
)
p = tf.add_paragraph()
run(
    p,
    "How it works: while the main loop runs, the kernel pre-fetches a ~24 KB window at the store path it is about to take.",
    13,
    SUB,
)

# LEFT panel: RUNTIME
box(0.45, 1.55, 6.15, 0.5, fill=PANEL2)
tf = box(0.45, 1.6, 6.15, 0.4)
run(tf.paragraphs[0], "RUNTIME  -  what picks the prefetch path?", 15, ACC, b=True)
rt = [
    ("GSU", GREEN, "1 -> normal store   |   >1 -> reduction store"),
    ("Beta", ACC, "beta = 0  vs  beta != 0  (two paths)"),
    ("Alpha", YELLOW, "no effect on the path (used inside it)"),
    ("Edge M/N", RED, "ragged size -> edge store may miss the window"),
]
y = 2.25
for name, c, desc in rt:
    dot(0.55, y + 0.05, c)
    tf = box(0.85, y - 0.05, 5.7, 0.55)
    p = tf.paragraphs[0]
    run(p, name + "  ", 14, c, b=True)
    run(p, desc, 12.5, FG)
    y += 0.62
tf = box(0.55, y + 0.02, 5.9, 0.9)
run(
    tf.paragraphs[0],
    "Test: sweep GSU x Beta x Alpha x edge; check the right path runs (arm-hit), entry latency, and numeric correctness.",
    12,
    SUB,
)

# RIGHT panel: BUILD-TIME
box(6.75, 1.55, 6.15, 0.5, fill=PANEL2)
tf = box(6.75, 1.6, 6.15, 0.4)
run(tf.paragraphs[0], "BUILD-TIME  -  which kernels are covered?", 15, ACC, b=True)
bt = [
    ("Pure GEMM (bbs / f8f8s / sss)", GREEN, "covered"),
    ("Sparse GEMM (spmm)", GREEN, "covered (same as dense)"),
    ("Fused activation", YELLOW, "covered, but activation body not"),
    ("MBSK / GSU>1 reduction", YELLOW, "writeback yes, reduction loop no"),
    ("StreamK", RED, "excluded -> uses PC-rel prefetch"),
    ("Subtile", RED, "StreamK:3 excluded; :0 unverified"),
]
y = 2.25
for name, c, verdict in bt:
    dot(6.85, y + 0.05, c)
    tf = box(7.15, y - 0.05, 5.6, 0.5)
    p = tf.paragraphs[0]
    run(p, name, 13, FG, b=True)
    p2 = tf.add_paragraph()
    run(p2, verdict, 11.5, c)
    y += 0.66

# Bottom: what needs work
box(0.45, 6.55, 12.45, 0.7, fill=RED)
tf = box(0.55, 6.62, 12.3, 0.6)
p = tf.paragraphs[0]
run(p, "Needs work:  ", 13, RGBColor(0x22, 0x14, 0x16), b=True)
run(
    p,
    "(1) large-tile activation function body   (2) GSU>1 reduction loop   -   both sit past the prefetch window.",
    13,
    RGBColor(0x22, 0x14, 0x16),
    b=True,
)

out = "/home/geotseng/workspace/fork/rocm-libraries/shared/stinkytofu/docs/abs_prefetch_experiment_plan.pptx"
prs.save(out)
print("saved", out, "slides=", len(prs.slides._sldIdLst))
