#!/usr/bin/env python3
"""Generate the abs-dynamic-prefetch experiment-design deck (gfx1250)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette (flat, minimal) ----
BG = RGBColor(0x14, 0x14, 0x14)
CARD = RGBColor(0x1E, 0x1E, 0x1E)
FG = RGBColor(0xE8, 0xE8, 0xE8)
SUB = RGBColor(0xA8, 0xA8, 0xA8)
ACC = RGBColor(0x59, 0x9C, 0xE7)  # blue accent
GREEN = RGBColor(0x3F, 0xA2, 0x66)
YELLOW = RGBColor(0xE0, 0xB0, 0x50)
RED = RGBColor(0xE0, 0x5A, 0x6E)
GRIDLN = RGBColor(0x3A, 0x3A, 0x3A)
HDRBG = RGBColor(0x2A, 0x36, 0x48)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def _tb(slide, l, t, w, h):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    return box, tf


def _run(p, text, size, color=FG, bold=False, italic=False, font="Calibri"):
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = font
    return r


def header(slide, kicker, title):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACC
    bar.line.fill.background()
    _, tf = _tb(slide, Inches(0.55), Inches(0.35), Inches(12.2), Inches(1.0))
    p = tf.paragraphs[0]
    _run(p, kicker.upper(), 12, ACC, bold=True)
    p2 = tf.add_paragraph()
    _run(p2, title, 28, FG, bold=True)


def title_slide():
    s = prs.slides.add_slide(BLANK)
    _bg(s)
    bar = s.shapes.add_shape(1, 0, Inches(2.7), SW, Inches(0.14))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACC
    bar.line.fill.background()
    _, tf = _tb(s, Inches(0.7), Inches(3.0), Inches(12), Inches(2.6))
    p = tf.paragraphs[0]
    _run(p, "gfx1250 abs 指令預取(dynamic)", 40, FG, bold=True)
    p2 = tf.add_paragraph()
    _run(p2, "有效性與支援度 — 實驗設計簡報", 26, ACC, bold=True)
    p3 = tf.add_paragraph()
    p3.space_before = Pt(16)
    _run(p3, "Runtime(launch 參數)× Build-time(kernel 類型)分軸驗證", 16, SUB)
    p4 = tf.add_paragraph()
    _run(
        p4,
        "來源:/data0/geotseng/comparison_output(2924 kernels)+ TensileLite/StinkyTofu 原始碼 · 6 個 read-only agent 交叉驗證",
        12,
        SUB,
        italic=True,
    )


def bullets_slide(kicker, title, items):
    s = prs.slides.add_slide(BLANK)
    _bg(s)
    header(s, kicker, title)
    _, tf = _tb(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.4))
    first = True
    for lvl, text, color in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(7)
        bullet = "▸ " if lvl == 0 else "– "
        _run(p, bullet, 16, ACC if lvl == 0 else SUB, bold=(lvl == 0))
        _run(p, text, 17 if lvl == 0 else 14.5, color, bold=(lvl == 0))
    return s


def table_slide(kicker, title, headers, rows, col_w, row_colors=None, note=None, fs=11):
    s = prs.slides.add_slide(BLANK)
    _bg(s)
    header(s, kicker, title)
    nrows, ncols = len(rows) + 1, len(headers)
    total_w = sum(col_w)
    left = Inches((13.333 - total_w) / 2)
    top = Inches(1.65)
    height = Inches(min(4.9, 0.42 * nrows))
    gtbl = s.shapes.add_table(nrows, ncols, left, top, Inches(total_w), height)
    tbl = gtbl.table
    tbl.first_row = False
    tbl.horz_banding = False
    for j, w in enumerate(col_w):
        tbl.columns[j].width = Inches(w)
    # header row
    for j, htext in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = HDRBG
        c.margin_left = Inches(0.06)
        c.margin_right = Inches(0.04)
        c.margin_top = Inches(0.02)
        c.margin_bottom = Inches(0.02)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = c.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _run(p, htext, fs + 0.5, FG, bold=True)
    # body
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = tbl.cell(i + 1, j)
            c.fill.solid()
            c.fill.fore_color.rgb = CARD if i % 2 == 0 else RGBColor(0x24, 0x24, 0x24)
            c.margin_left = Inches(0.06)
            c.margin_right = Inches(0.04)
            c.margin_top = Inches(0.02)
            c.margin_bottom = Inches(0.02)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = c.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            color = FG
            if row_colors and j == row_colors[0]:
                color = row_colors[1][i]
            _run(p, str(val), fs, color, bold=(j == 0))
    if note:
        _, tf = _tb(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.7))
        p = tf.paragraphs[0]
        _run(p, note, 11, SUB, italic=True)
    return s


# ============================ SLIDES ============================
title_slide()

bullets_slide(
    "01 · 背景",
    "目標與方法",
    [
        (
            0,
            "評估 gfx1250『abs 指令預取(dynamic)』對 post-loop / Global-Write(GW)epilogue 的覆蓋有效性與各 kernel 類型支援度",
            FG,
        ),
        (0, "分兩軸設計實驗", ACC),
        (
            1,
            "Runtime:固定 kernel binary,變動 launch 參數(GSU / Alpha / Beta / Edge M,N)—— 這些是 abs dynamic 選臂條件",
            FG,
        ),
        (
            1,
            "Build-time:kernel 類型是否被 abs 支援(pure-GEMM / fused / MBSK / StreamK / sparse / subtile)",
            FG,
        ),
        (0, "有效性指標", ACC),
        (1, "Arm-hit:runtime 選到的 epilogue 是否命中 3 臂之一(覆蓋正確性)", FG),
        (1, "Epilogue 進入延遲 / kernel 時間;數值正確性回歸", FG),
        (
            0,
            "本簡報為實驗設計 + 支援矩陣 + 依覆蓋分析推導的預期結果(尚無 perf 實測)",
            SUB,
        ),
    ],
)

table_slide(
    "02 · 機制",
    "三段 regime(由 totalLayoutBytes 決定)",
    ["regime", "totalLayoutBytes", "負責 pass", "post-loop 覆蓋方式", "kernel 數"],
    [
        [
            "入口 preload",
            "≤ 32640",
            "dynamic no-op",
            "CP 啟動常駐 [0,32640),整顆全覆蓋",
            "1622",
        ],
        [
            "static ladder",
            "(32640, 65536]",
            "AbsStaticPass",
            "每 4096B 線性佈點,GW 各段連續填滿",
            "649",
        ],
        [
            "dynamic 3-arm",
            "> 65536",
            "AbsDynamicPass",
            "MultiGemmEnd 後 1 組 3-arm 定向 burst",
            "653",
        ],
    ],
    [1.9, 2.2, 2.0, 4.5, 1.1],
    row_colors=(0, [GREEN, ACC, YELLOW]),
    note="重點:大 kernel 的 dynamic 檔開頭每-4096B『planned insert sites』只是 debug 預覽、不 emit;>65536B 實際只發 3-arm burst(每臂 N=6×4096=24576B)。",
)

bullets_slide(
    "03 · 機制",
    "D1 動態 3-arm ladder(emit 於 label_MultiGemmEnd 後)",
    [
        (0, "執行期以 masked GSU 與 Beta 重新判斷,選一條臂 prefetch(每臂 24576B)", FG),
        (1, "Arm A:GSU > 1  →  GW_B0_MB / GW_B0_MBSK(workspace reduction 寫回)", GREEN),
        (1, "Arm B:GSU == 1 & Beta == 0  →  GW_B0_GSU1(單 GSU 直寫)", ACC),
        (1, "Arm C:GSU == 1 & Beta ≠ 0  →  GW_B1_GSU1(單 GSU + βC,DEFAULT)", YELLOW),
        (0, "臂數對應執行期『互斥』的 GSU/beta regime(≤3)→ 無 #variants 缺口", FG),
        (
            0,
            "MB 臂為 B0-only(GSU>1 MultipleBuffer 時 hasBeta 強制 false;β 在 reduction kernel 套用)",
            SUB,
        ),
        (
            0,
            "程式碼:KernelWriterAssembly.py:15006–15012(GSU)、14562–14585(Beta);C++ ladder:SwInstructionPrefetchAbsDynamicPass.cpp:504–536",
            SUB,
        ),
    ],
)

# Runtime section divider-ish
bullets_slide(
    "04 · Runtime 實驗",
    "launch 參數 → 選臂映射(原始碼驗證)",
    [
        (0, "GSU(masked)= 選 A vs {B,C} 的唯一開關", GREEN),
        (1, "GSU∈{2,4,8} → Arm A;GSU==1 → 進 B/C  (KWA.py:15006–15012)", FG),
        (0, "Beta = 在 GSU==1 下選 B vs C(純 zero / non-zero 測試)", ACC),
        (1, "β==0 → B0_GSU1;β≠0 → B1_GSU1;無 β==1 特例  (KWA.py:14562–14585)", FG),
        (0, "Alpha ✗ 不改變選臂", YELLOW),
        (
            1,
            "無 s[Alpha] 分支;alpha 於被選 block 內以 v_mul 套用(GlobalWriteBatch.py:2866+)。改變的是窗內執行指令,非進入點",
            FG,
        ),
        (0, "Edge(M/N)= per-workgroup tile 位置的 runtime 決策", RED),
        (
            1,
            "僅『邊界 WG 且 Size%MT≠0』走 Edge/VW1;內部 WG 走 NonEdge/VW8  (KWA.py:14817–14876)",
            FG,
        ),
        (
            1,
            "Edge/VW1 子塊常在 24576B 窗『之外』→ 覆蓋取決於 layout 距離,無 edge-aware 重錨",
            RED,
        ),
    ],
)

table_slide(
    "05 · Runtime 實驗",
    "測試矩陣(假設 >65536B,DYNAMIC regime)",
    ["#", "GSU", "Beta", "Alpha", "M/N", "選中臂", "覆蓋?", "斷言"],
    [
        [
            "1",
            "2",
            "n/a",
            "1",
            "aligned",
            "A GW_B0_MB",
            "是",
            "arm-hit A · 進入延遲 · D=αAB",
        ],
        [
            "3",
            "8",
            "0",
            "1",
            "+1",
            "A + edge",
            "入口是/VW1外",
            "邊界 WG edge 未覆蓋 · 尾端正確性",
        ],
        ["4", "1", "0", "1", "aligned", "B GW_B0_GSU1", "是", "arm-hit B · 進入延遲"],
        ["6", "1", "1", "1", "aligned", "C GW_B1_GSU1", "是", "arm-hit C · D=αAB+βC"],
        [
            "7",
            "1",
            "k≠0",
            "k",
            "aligned",
            "C GW_B1_GSU1",
            "是",
            "β≠0,≠1 同臂(無 β==1 特例)",
        ],
        [
            "8",
            "1",
            "k≠0",
            "1",
            "+1",
            "C + edge",
            "入口是/VW1外",
            "邊界 edge 未覆蓋 · C-load 正確性",
        ],
        [
            "9",
            "1",
            "0",
            "0",
            "aligned",
            "B GW_B0_GSU1",
            "是",
            "α=0 不改變臂;窗內結果歸零",
        ],
        [
            "10",
            "2",
            "1",
            "1",
            "aligned",
            "A GW_B0_MB",
            "是",
            "GSU>1 時 beta 不影響選臂",
        ],
    ],
    [0.5, 0.8, 0.9, 0.9, 1.1, 2.2, 1.5, 4.2],
    row_colors=(5, [GREEN, GREEN, ACC, YELLOW, YELLOW, YELLOW, ACC, GREEN]),
    fs=10.5,
    note="獨立性:單獨變 GSO(1/4/6 行)隔離 GSU;GSU=1 下變 Beta(4 vs 6)隔離 beta;固定 GSU/Beta 變 Alpha(4/9、6/7)證明 alpha 不移臂。",
)

bullets_slide(
    "06 · Runtime 實驗",
    "有效性量測方法",
    [
        (0, "Arm-hit(覆蓋正確性)", GREEN),
        (
            1,
            "驗證 runtime 執行到的 GW_B* label 正是 GSU/Beta 比較所選的臂;比對 3-arm ladder 的 prefetch 目標",
            FG,
        ),
        (
            1,
            "『covered』= 執行的 store 碼距臂入口 byte offset < 24576;Edge/VW1 是最可能落窗外者",
            FG,
        ),
        (0, "Epilogue 進入延遲 / kernel 時間", ACC),
        (
            1,
            "同一 binary、同尺寸,比較 abs on vs PC-rel:量 epilogue 進入點的前端 stall / 整體 kernel 時間",
            FG,
        ),
        (0, "數值正確性回歸", YELLOW),
        (
            1,
            "掃 GSU×Beta×Alpha×Edge 組合,驗證 D 結果正確(α=0、β=0、ragged tail 邊界皆需正確)",
            FG,
        ),
        (0, "建議:先做 arm-hit + 正確性(無需硬體計數器),再視情況補延遲量測", SUB),
    ],
)

table_slide(
    "07 · Build-time",
    "kernel 類型支援矩陣",
    [
        "kernel 類型",
        "abs base 配置?",
        "典型 regime",
        "GW-tree 覆蓋?",
        "已知盲區",
        "fallback",
    ],
    [
        [
            "pure-GEMM (bbs/f8f8s/sss)",
            "是",
            "全三段",
            "是(3-case)",
            "OptNLL 快路徑(>32640)",
            "—",
        ],
        [
            "fused activation",
            "是",
            "小=入口/大=dyn",
            "是",
            "activation 函式本體(D2+)",
            "—",
        ],
        [
            "MBSK / GSU>1 reduction",
            "是",
            "dynamic",
            "寫回塊是",
            "reduction loop 本體(D2+)",
            "—",
        ],
        ["StreamK", "否(-1)", "—", "N/A", "全 kernel", "PC-rel"],
        [
            "sparse (spmm)",
            "是",
            "幾乎全 dyn",
            "是(同 dense)",
            "同 dense + tail_coalesced",
            "—",
        ],
        ["subtile", "條件式", "見下頁", "部分/未驗證", "StreamK:3 時全排除", "PC-rel*"],
    ],
    [3.0, 1.5, 1.9, 1.8, 2.6, 1.0],
    row_colors=(5, [GREEN, YELLOW, YELLOW, ACC, GREEN, RED]),
    fs=10.5,
    note='guard:KernelWriter.py:9623 `if SwInstructionPrefetchAbs and not kernel["StreamK"] and version==(12,5,0)` → StreamK 使 baseSgpr=-1 → abs 全 no-op。',
)

bullets_slide(
    "08 · Build-time",
    "subtile 特別說明(UseSubtileImpl ≠ StreamK)",
    [
        (0, "subtile 的支援度取決於該解的 StreamK 設定", YELLOW),
        (
            1,
            "subtile_bf16_gfx1250.yaml(主功能測試)= StreamK:[3] → 被 line-9623 guard 排除 → baseSgpr=-1 → abs no-op → 回退 PC-rel(等同 StreamK)",
            RED,
        ),
        (
            1,
            "bench / cluster 設定 = StreamK:[0](GSU 強制 1)→ 會配置 abs base,理論上進 static/dynamic",
            GREEN,
        ),
        (
            0,
            "但 GSU1-subtile 路徑的 GW-tree(dynamic)與 entry-burst(static)覆蓋『尚未 fleet 驗證』",
            FG,
        ),
        (
            0,
            "→ 建議實驗:針對 StreamK:0 的 subtile 解,驗證 arm-hit 與 subtile edge(checkIsEdgeSubtile, KWA.py:14602)是否落窗內",
            SUB,
        ),
    ],
)

table_slide(
    "09 · 盲區",
    "4 個覆蓋盲區(agent 驗證)",
    ["#", "盲區", "位置", "影響範疇", "嚴重度"],
    [
        [
            "G1",
            "共用 activation 函式本體",
            "kernel 尾端(s_swappc 目標),bf16 217068–221200",
            "大型 fused / sparse",
            "高",
        ],
        [
            "G2",
            "GSU>1 reduction 累加 loop",
            "workspace 讀-累加(僅寫回塊被 arm A 覆蓋)",
            "MBSK / GSU>1",
            "中",
        ],
        [
            "G3",
            "OptNLL 快速路徑 store",
            "GW_B0_OptNLL_MB 非 D1 臂,>32640B 即盲",
            "所有含 OptNLL 大 kernel",
            "中",
        ],
        [
            "G4",
            "巨型 activation-dispatch tail_coalesced",
            "18–24KB block 只 1 prefetch,尾端 ~16–20KB 無新預取",
            "12-way activation",
            "中",
        ],
    ],
    [0.5, 3.0, 4.6, 2.6, 1.0],
    row_colors=(4, [RED, YELLOW, YELLOW, YELLOW]),
    fs=10.5,
    note="C++ D2+ deferred 標記:header :37/:126;detector :317/:320-321/:322/:327-328/:333;bail :381-385。普通 3-case GW 樹本身無缺口。",
)

bullets_slide(
    "10 · 結論",
    "結論與建議",
    [
        (
            0,
            "abs dynamic 對 pure-GEMM / sparse 的 GW dispatch 樹已完整覆蓋;3 臂對應互斥 GSU/beta regime,無 #variants 缺口",
            GREEN,
        ),
        (0, "真正需要處理的盲區(與初始判斷一致)", RED),
        (1, "大 tile activation 函式本體(G1/G4)、GSU>1 reduction loop 本體(G2)", FG),
        (0, "天生排除、走 PC-relative:StreamK 與 StreamK:3 的 subtile", ACC),
        (0, "建議下一步", ACC),
        (
            1,
            "Runtime:先跑 arm-hit + 正確性掃描(GSU×β×α×edge),確認 3 臂命中與 ragged-tail 正確",
            FG,
        ),
        (
            1,
            "Build-time:驗證 StreamK:0 subtile 的覆蓋;評估把 activation 本體 / reduction loop 納入 D2+ 預取或增大 N",
            FG,
        ),
        (
            0,
            "潛在修法:對 s_swappc activation 目標另發預取,或把本體移進某臂窗內 / 增大該臂 N",
            SUB,
        ),
    ],
)

out = "/data0/geotseng/abs_prefetch_experiment_plan.pptx"
prs.save(out)
print("saved", out, "slides=", len(prs.slides._sldIdLst))
