#!/usr/bin/env python3
"""Single-slide abs-dynamic-prefetch experiment summary (gfx1250)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR

BG = RGBColor(0x14, 0x14, 0x14)
CARD = RGBColor(0x1E, 0x1E, 0x1E)
CARD2 = RGBColor(0x24, 0x24, 0x24)
FG = RGBColor(0xE8, 0xE8, 0xE8)
SUB = RGBColor(0xA8, 0xA8, 0xA8)
ACC = RGBColor(0x59, 0x9C, 0xE7)
GREEN = RGBColor(0x3F, 0xA2, 0x66)
YELLOW = RGBColor(0xE0, 0xB0, 0x50)
RED = RGBColor(0xE0, 0x5A, 0x6E)
HDR = RGBColor(0x2A, 0x36, 0x48)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid()
s.background.fill.fore_color.rgb = BG
SW = prs.slide_width


def run(p, t, sz, c=FG, b=False, i=False):
    r = p.add_run()
    r.text = t
    f = r.font
    f.size = Pt(sz)
    f.bold = b
    f.italic = i
    f.color.rgb = c
    f.name = "Calibri"
    return r


def box(l, t, w, h):
    bx = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = bx.text_frame
    tf.word_wrap = True
    return tf


def tbl(l, t, w, headers, rows, colw, fs, hcolor_col=None, hcolors=None, rowh=0.3):
    n = len(rows) + 1
    g = s.shapes.add_table(
        n, len(headers), Inches(l), Inches(t), Inches(w), Inches(rowh * n)
    ).table
    g.first_row = False
    g.horz_banding = False
    for j, cw in enumerate(colw):
        g.columns[j].width = Inches(cw)
    for j, h in enumerate(headers):
        c = g.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = HDR
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(c, m, Inches(0.03))
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        run(c.text_frame.paragraphs[0], h, fs, FG, b=True)
    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            c = g.cell(i + 1, j)
            c.fill.solid()
            c.fill.fore_color.rgb = CARD if i % 2 == 0 else CARD2
            for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
                setattr(c, m, Inches(0.03))
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            col = FG
            if hcolor_col is not None and j == hcolor_col and hcolors:
                col = hcolors[i]
            run(c.text_frame.paragraphs[0], str(v), fs, col, b=(j == 0))


# ---- top accent + title ----
bar = s.shapes.add_shape(1, 0, 0, SW, Inches(0.1))
bar.fill.solid()
bar.fill.fore_color.rgb = ACC
bar.line.fill.background()
tf = box(0.4, 0.18, 12.6, 0.9)
p = tf.paragraphs[0]
run(p, "gfx1250 abs 指令預取(dynamic)— 有效性與支援度實驗", 22, FG, b=True)
p2 = tf.add_paragraph()
run(
    p2,
    "Regime(依 totalLayoutBytes):入口 preload ≤32640 (n=1622) · static ladder (32640,65536] (n=649) · dynamic 3-arm >65536 (n=653) — 每臂 N=6×4096=24576B",
    10.5,
    SUB,
)

# ---- LEFT: Runtime ----
tf = box(0.4, 1.35, 6.3, 0.4)
run(tf.paragraphs[0], "RUNTIME(固定 binary,變 launch 參數 → 選臂)", 13, ACC, b=True)
tbl(
    0.4,
    1.75,
    6.35,
    ["參數", "作用 / 選臂", "原始碼"],
    [
        ["GSU", "masked;>1→A, ==1→B/C(唯一選 A 開關)", "KWA:15006"],
        ["Beta", "==0→B0, ≠0→B1;無 β==1 特例", "KWA:14562"],
        ["Alpha", "✗ 不選臂;窗內 v_mul 套用", "GWB:2866"],
        ["Edge M/N", "邊界WG且Size%MT≠0→VW1(常在窗外)", "KWA:14817"],
    ],
    [1.0, 4.15, 1.2],
    9.5,
    hcolor_col=0,
    hcolors=[GREEN, ACC, YELLOW, RED],
    rowh=0.4,
)

tf = box(0.4, 3.75, 6.3, 0.35)
run(
    tf.paragraphs[0],
    "測試矩陣(DYNAMIC regime)· 指標:arm-hit / 進入延遲 / 正確性",
    11,
    ACC,
    b=True,
)
tbl(
    0.4,
    4.1,
    6.35,
    ["GSU", "β", "α", "M/N", "臂", "覆蓋"],
    [
        ["2", "-", "1", "aln", "A B0_MB", "是"],
        ["1", "0", "1", "aln", "B B0_GSU1", "是"],
        ["1", "1", "1", "aln", "C B1_GSU1", "是"],
        ["1", "k", "1", "+1", "C+edge", "入口是/VW1外"],
        ["1", "0", "0", "aln", "B B0_GSU1", "α=0 不移臂"],
    ],
    [0.7, 0.6, 0.6, 0.7, 2.0, 1.75],
    9.5,
    hcolor_col=4,
    hcolors=[GREEN, ACC, YELLOW, YELLOW, ACC],
    rowh=0.34,
)

# ---- RIGHT: Build-time ----
tf = box(6.95, 1.35, 6.0, 0.4)
run(tf.paragraphs[0], "BUILD-TIME(kernel 類型支援度)", 13, ACC, b=True)
tbl(
    6.95,
    1.75,
    6.0,
    ["kernel 類型", "abs?", "已知盲區"],
    [
        ["pure-GEMM (bbs/f8f8s/sss)", "是", "OptNLL 快路徑(>32640)"],
        ["fused activation", "是", "activation 本體 (D2+)"],
        ["MBSK / GSU>1 reduction", "是", "reduction loop 本體 (D2+)"],
        ["sparse (spmm)", "是", "同 dense + tail_coalesced"],
        ["StreamK", "否→PC-rel", "全 kernel"],
        ["subtile", "條件式", "StreamK:3 排除;:0 未驗證"],
    ],
    [3.1, 1.05, 1.85],
    9.5,
    hcolor_col=1,
    hcolors=[GREEN, YELLOW, YELLOW, GREEN, RED, RED],
    rowh=0.38,
)

tf = box(6.95, 4.35, 6.0, 0.35)
run(tf.paragraphs[0], "覆蓋盲區(agent 驗證;3-case GW 樹本身無缺口)", 11, ACC, b=True)
tf = box(6.95, 4.7, 6.0, 2.0)
for gid, txt, c in [
    ("G1", "activation 函式本體(尾端 s_swappc)— 大 fused/sparse", RED),
    ("G2", "GSU>1 reduction 累加 loop 本體 — 僅寫回塊被覆蓋", YELLOW),
    ("G3", "OptNLL 快路徑 store(非 D1 臂,>32640 即盲)", YELLOW),
    ("G4", "巨型 activation-dispatch tail_coalesced 尾端 ~16–20KB", YELLOW),
]:
    p = tf.add_paragraph()
    p.space_after = Pt(3)
    run(p, gid + "  ", 10.5, c, b=True)
    run(p, txt, 10.5, FG)

# ---- bottom conclusion strip ----
bar2 = s.shapes.add_shape(1, Inches(0.4), Inches(6.75), Inches(12.55), Inches(0.55))
bar2.fill.solid()
bar2.fill.fore_color.rgb = HDR
bar2.line.fill.background()
tf = bar2.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.12)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
run(p, "結論:", 11, ACC, b=True)
run(
    p,
    "pure/sparse GW 樹已完整覆蓋、3 臂互斥無 #variants 缺口;需解 = 大 tile activation 本體(G1/G4)+ GSU>1 reduction loop(G2);StreamK 與 StreamK:3 subtile 走 PC-rel。",
    10.5,
    FG,
)

out = "/data0/geotseng/abs_prefetch_experiment_plan.pptx"
prs.save(out)
print("saved", out, "slides=", len(prs.slides._sldIdLst))
